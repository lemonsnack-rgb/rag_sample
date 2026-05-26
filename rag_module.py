import os
import re
import io
import csv
import hashlib
import time
from datetime import datetime
import streamlit as st

# 🔧 긴급 수정: .env 파일 강제 로드
from dotenv import load_dotenv
load_dotenv(override=True)  # 기존 환경변수 덮어쓰기

# Google Drive API import
import google.auth
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# Supabase Imports
from supabase import create_client, Client

# LangChain and Embedding Imports
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_google_genai import GoogleGenerativeAIEmbeddings

# File Parsing Libraries Imports
import pypdf
import docx
import openpyxl
from pptx import Presentation
from PIL import Image
import pytesseract

EMBEDDING_MODEL = "models/gemini-embedding-001"
EMBEDDING_DIMENSION = 768
EMBEDDING_BATCH_SIZE = int(os.getenv("EMBEDDING_BATCH_SIZE", "4"))
EMBEDDING_MAX_RETRIES = int(os.getenv("EMBEDDING_MAX_RETRIES", "5"))
EMBEDDING_RETRY_BASE_SECONDS = float(os.getenv("EMBEDDING_RETRY_BASE_SECONDS", "8"))
EMBEDDING_BATCH_DELAY_SECONDS = float(os.getenv("EMBEDDING_BATCH_DELAY_SECONDS", "2"))


def create_embeddings():
    return GoogleGenerativeAIEmbeddings(
        model=EMBEDDING_MODEL,
        output_dimensionality=EMBEDDING_DIMENSION,
    )


def embed_documents_with_retry(embeddings, texts):
    embedding_vectors = []
    batch_size = max(1, EMBEDDING_BATCH_SIZE)

    for start in range(0, len(texts), batch_size):
        batch = texts[start:start + batch_size]
        batch_no = (start // batch_size) + 1
        total_batches = (len(texts) + batch_size - 1) // batch_size

        for attempt in range(EMBEDDING_MAX_RETRIES):
            try:
                embedding_vectors.extend(embeddings.embed_documents(batch))
                break
            except Exception as e:
                if attempt >= EMBEDDING_MAX_RETRIES - 1:
                    raise

                wait_seconds = EMBEDDING_RETRY_BASE_SECONDS * (2 ** attempt)
                message = str(e)
                if "RESOURCE_EXHAUSTED" not in message and "429" not in message:
                    wait_seconds = min(wait_seconds, 10)

                st.warning(
                    f"임베딩 API 제한으로 대기 중입니다. "
                    f"배치 {batch_no}/{total_batches}, 재시도 {attempt + 1}/{EMBEDDING_MAX_RETRIES}, "
                    f"{wait_seconds:.0f}초 후 재시도합니다."
                )
                time.sleep(wait_seconds)

        if start + batch_size < len(texts) and EMBEDDING_BATCH_DELAY_SECONDS > 0:
            time.sleep(EMBEDDING_BATCH_DELAY_SECONDS)

    return embedding_vectors

# ==================== [텍스트 전처리 - 메타데이터 분리 개선] ====================
def preprocess_text_with_section_headers(text):
    """
    섹션 태그를 메타데이터로 분리하여 임베딩 품질 향상
    """
    if text:
        text = text.replace('\x00', '')

    lines = text.split('\n')
    chunks = []
    current_section = "일반"
    header_pattern = re.compile(r'^\s*제\s*\d+\s*(조|장)')

    current_chunk_lines = []

    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue

        # 이미 태그가 있는 경우 (Excel, PPT 등)
        if stripped_line.startswith('[') and ']' in stripped_line:
            if current_chunk_lines:
                chunks.append({
                    "content": "\n".join(current_chunk_lines),
                    "section": current_section
                })
                current_chunk_lines = []
            # 태그 제거하고 순수 내용만 저장
            content_only = re.sub(r'^\[.*?\]\s*', '', stripped_line)
            current_chunk_lines.append(content_only)
            continue

        if header_pattern.match(stripped_line):
            # 새로운 섹션 시작
            if current_chunk_lines:
                chunks.append({
                    "content": "\n".join(current_chunk_lines),
                    "section": current_section
                })
                current_chunk_lines = []

            # 🔧 개선: 섹션명을 첫 100자로 제한 (긴 제목 방지)
            current_section = stripped_line[:100] if len(stripped_line) > 100 else stripped_line
            current_chunk_lines.append(stripped_line)
        else:
            # 순수 내용만 저장 (태그 없이)
            current_chunk_lines.append(stripped_line)

    # 마지막 청크 추가
    if current_chunk_lines:
        chunks.append({
            "content": "\n".join(current_chunk_lines),
            "section": current_section
        })

    return chunks

# ==================== [파일 포맷별 텍스트 추출 (OCR 강화)] ====================
def extract_text_from_pdf(fh):
    text = ""
    try:
        reader = pypdf.PdfReader(fh)
        for page in reader.pages:
            t = page.extract_text()
            if t: text += t + "\n"
            else:
                try:
                    for image_file in page.images:
                        img = Image.open(io.BytesIO(image_file.data))
                        text += pytesseract.image_to_string(img, lang='kor+eng') + "\n"
                except: pass
    except Exception as e: print(f"PDF Error: {e}")
    return text

def extract_text_from_docx(fh):
    text = ""
    try:
        doc = docx.Document(fh)
        for para in doc.paragraphs: text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                row_text = [c.text.strip() for c in row.cells]
                text += " | ".join(row_text) + "\n"
            text += "\n"
    except Exception as e: print(f"DOCX Error: {e}")
    return text

def extract_text_from_xlsx(fh, fname):
    text = ""
    try:
        wb = openpyxl.load_workbook(fh, data_only=True)
        for sname in wb.sheetnames:
            sheet = wb[sname]
            rows = list(sheet.rows)
            if not rows: continue
            headers = [str(c.value).strip() if c.value else f"열{i}" for i,c in enumerate(rows[0])]
            for row in rows[1:]:
                parts = []
                for i, c in enumerate(row):
                    if i < len(headers) and c.value is not None:
                        val = str(c.value).strip()
                        if val: parts.append(f"{headers[i]}: {val}")
                if parts: text += f"{fname}-{sname} " + ", ".join(parts) + "\n"
    except Exception as e: print(f"XLSX Error: {e}")
    return text

def extract_text_from_pptx(fh):
    text = ""
    try:
        prs = Presentation(fh)
        for i, slide in enumerate(prs.slides):
            stext = []
            for shape in slide.shapes:
                if hasattr(shape, "text"): stext.append(shape.text)
            if stext: text += f"슬라이드 {i+1} " + "\n".join(stext) + "\n"
    except Exception as e: print(f"PPTX Error: {e}")
    return text

def extract_text_from_txt(fh):
    try:
        c = fh.read()
        try: return c.decode('utf-8')
        except: return c.decode('cp949')
    except: return ""

def extract_text_from_csv(fh, fname):
    text = ""
    try:
        c = fh.read()
        try: dec = c.decode('utf-8')
        except: dec = c.decode('cp949')
        rows = list(csv.reader(io.StringIO(dec)))
        if not rows: return ""
        headers = rows[0]
        for row in rows[1:]:
            parts = []
            for i, v in enumerate(row):
                if i < len(headers) and v.strip():
                    parts.append(f"{headers[i]}: {v.strip()}")
            if parts: text += f"{fname} " + ", ".join(parts) + "\n"
    except Exception as e: print(f"CSV Error: {e}")
    return text

def extract_text_from_image(fh):
    text = ""
    try:
        img = Image.open(fh)
        text = pytesseract.image_to_string(img, lang='kor+eng')
    except Exception as e: print(f"OCR Error: {e}")
    return text

# ==================== [파일 포맷별 최적 청크 크기] ====================
def get_optimal_splitter(file_type):
    """
    파일 타입에 따라 최적화된 청크 크기 반환

    최적화: 청크 크기 축소로 키워드 밀도 증가
    - 800자: 키워드가 청크 내에서 높은 비중 차지
    - 정확한 매칭 가능성 증가
    - 규정집처럼 조항 단위가 중요한 문서에 최적
    """
    if file_type == 'xlsx':
        return RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    elif file_type == 'pptx':
        return RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    elif file_type == 'csv':
        return RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    else:
        # PDF, DOCX 등 일반 문서: 800자 권장
        return RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)

# ==================== [핵심 로직: Supabase 연결 및 동기화] ====================

def init_vector_store():
    """
    Supabase 클라이언트를 생성하고 임베딩 모델을 초기화합니다.
    """
    url = os.environ.get("SUPABASE_URL")
    key = os.environ.get("SUPABASE_KEY")

    if not url or not key:
        raise ValueError("SUPABASE_URL 또는 SUPABASE_KEY 환경 변수가 설정되지 않았습니다.")

    try:
        supabase_client: Client = create_client(url, key)
        embeddings = create_embeddings()

        return {
            'supabase_client': supabase_client,
            'embeddings': embeddings
        }
    except Exception as e:
        print(f"Supabase 클라이언트 초기화 실패: {e}")
        raise

def get_file_hash(content):
    """파일 내용의 해시값 생성 (중복 체크용)"""
    return hashlib.md5(content.encode('utf-8')).hexdigest()

def delete_document_by_source(client, source_name):
    """
    특정 파일명의 모든 문서를 DB에서 삭제
    중복 방지 및 업데이트를 위한 함수
    """
    try:
        result = client.table("documents").delete().eq("metadata->>source", source_name).execute()
        print(f"[OK] {source_name} 기존 데이터 삭제 완료")
        return True
    except Exception as e:
        print(f"[ERROR] {source_name} 삭제 실패: {e}")
        return False

def get_file_timestamps_from_db(supabase_client):
    """
    DB에 저장된 파일별 마지막 수정 시간을 조회합니다.

    Returns:
        dict: {filename: last_modified_timestamp}
    """
    try:
        # documents 테이블에서 파일별 최신 last_modified 조회
        result = supabase_client.table("documents").select("metadata").execute()

        file_times = {}
        for doc in result.data:
            metadata = doc.get('metadata', {})
            source = metadata.get('source')
            last_modified = metadata.get('last_modified')

            if source and last_modified:
                # 같은 파일의 여러 청크 중 가장 최신 시간 유지
                if source not in file_times or last_modified > file_times[source]:
                    file_times[source] = last_modified

        return file_times
    except Exception as e:
        print(f"DB 타임스탬프 조회 실패: {e}")
        return {}

def sync_drive_to_db(folder_id, supabase_client, force_update=False):
    """
    [개선된 핵심 함수] Google Drive에서 파일을 가져와 텍스트를 추출하고 벡터 DB에 동기화합니다.

    개선사항:
    - 🆕 증분 동기화: 변경된 파일만 자동 감지하여 업데이트
    - 중복 방지: 파일별로 기존 데이터 삭제 후 재삽입
    - 파일 타입별 최적 청크 크기
    - 섹션 태그를 메타데이터로 분리
    - 진행상황 상세 로그

    Args:
        folder_id: Google Drive 폴더 ID
        supabase_client: Supabase 클라이언트
        force_update: True일 경우 전체 재색인 (기본 False = 증분 동기화)
    """
    creds, _ = google.auth.default()
    service = build('drive', 'v3', credentials=creds)

    # Drive에서 파일 목록 가져오기 (modifiedTime 포함)
    res = service.files().list(
        q=f"'{folder_id}' in parents and trashed=false",
        fields="files(id, name, modifiedTime)"
    ).execute()
    files = res.get('files', [])

    st.write(f"🔍 폴더 내 파일 {len(files)}개 감지됨")

    # 증분 동기화: 변경된 파일만 감지
    files_to_process = []
    files_to_delete = []

    if not force_update:
        # DB에서 기존 파일의 last_modified 조회
        db_file_times = get_file_timestamps_from_db(supabase_client)
        drive_file_names = {f['name'] for f in files}

        new_count = 0
        updated_count = 0
        unchanged_count = 0

        for f in files:
            fname = f['name']
            drive_modified = f.get('modifiedTime', '')

            if fname not in db_file_times:
                # 새 파일
                files_to_process.append(f)
                new_count += 1
            elif drive_modified > db_file_times[fname]:
                # 수정된 파일
                files_to_process.append(f)
                updated_count += 1
            else:
                # 변경 없음
                unchanged_count += 1

        # Drive에 없지만 DB에 있는 파일 = 삭제된 파일
        for fname in db_file_times:
            if fname not in drive_file_names:
                files_to_delete.append(fname)

        # 변경 사항 요약
        st.info(f"""
        📊 증분 동기화 분석
        - 🆕 새 파일: {new_count}개
        - 🔄 수정된 파일: {updated_count}개
        - ✅ 변경 없음: {unchanged_count}개
        - 🗑️ 삭제된 파일: {len(files_to_delete)}개
        """)

        if len(files_to_process) == 0 and len(files_to_delete) == 0:
            st.success("✨ 모든 문서가 최신 상태입니다!")
            return 0
    else:
        # 전체 재색인 모드
        st.info("🔄 전체 재색인 모드 (모든 파일 처리)")
        files_to_process = files

    # 삭제된 파일 처리
    deleted_count = 0
    if files_to_delete:
        st.write("[DELETE] 삭제된 파일 정리 중...")
        for fname in files_to_delete:
            if delete_document_by_source(supabase_client, fname):
                deleted_count += 1
                st.caption(f"  [OK] {fname} 제거됨")

    # 임베딩 모델 초기화
    embeddings = create_embeddings()

    cnt = 0
    skipped = 0
    failed = 0
    total_to_process = len(files_to_process)

    if total_to_process == 0:
        st.success(f"[OK] 동기화 완료 (삭제: {deleted_count}개)")
        return deleted_count

    progress = st.progress(0)

    for i, f in enumerate(files_to_process):
        fid, fname = f['id'], f['name']
        drive_modified = f.get('modifiedTime', '')
        ext = fname.split('.')[-1].lower() if '.' in fname else ""
        progress.progress((i+1)/total_to_process, text=f"처리 중: {fname}")

        if ext not in ['pdf', 'docx', 'xlsx', 'pptx', 'txt', 'csv', 'md', 'jpg', 'jpeg', 'png']:
            st.caption(f"[SKIP] {fname} - 지원하지 않는 형식")
            skipped += 1
            continue

        try:
            # 파일 다운로드
            req = service.files().get_media(fileId=fid)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, req)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            fh.seek(0)

            # 텍스트 추출
            content = ""
            if ext == 'pdf':
                content = extract_text_from_pdf(fh)
            elif ext == 'docx':
                content = extract_text_from_docx(fh)
            elif ext == 'xlsx':
                content = extract_text_from_xlsx(fh, fname)
            elif ext == 'pptx':
                content = extract_text_from_pptx(fh)
            elif ext in ['txt', 'md']:
                content = extract_text_from_txt(fh)
            elif ext == 'csv':
                content = extract_text_from_csv(fh, fname)
            elif ext in ['jpg', 'png', 'jpeg']:
                content = extract_text_from_image(fh)

            if not content.strip():
                st.warning(f"[WARNING] {fname} - 내용 없음")
                skipped += 1
                continue

            # 섹션 인식 전처리 (메타데이터 분리)
            processed_chunks = preprocess_text_with_section_headers(content)

            # 파일 타입별 최적 청크 크기
            splitter = get_optimal_splitter(ext)

            # 문서 생성 (섹션 정보를 메타데이터 + 본문에 포함)
            docs = []
            for chunk_data in processed_chunks:
                # 추가 청크 분할
                sub_chunks = splitter.split_text(chunk_data["content"])
                for sub_chunk in sub_chunks:
                    if sub_chunk.strip():
                        section_name = chunk_data["section"]

                        # 🔧 개선: 섹션 태그 제거 - 순수 내용만 임베딩
                        # 이유:
                        # 1. 임베딩 토큰 100% 활용 (섹션 태그 오버헤드 제거)
                        # 2. 검색 노이즈 감소
                        # 3. "제5조" 등은 본문에 이미 포함되어 있어 검색 가능
                        docs.append(Document(
                            page_content=sub_chunk,  # 순수 내용만 (섹션 태그 없음)
                            metadata={
                                "source": fname,
                                "section": section_name,  # 메타데이터에는 유지 (참고용)
                                "file_type": ext,
                                "last_modified": drive_modified,
                                "created_at": datetime.now().isoformat()
                            }
                        ))

            if docs:
                # ✅ RPC 직접 호출 + embed_documents() (VECTOR(768) 저장)
                try:
                    # 문서용 임베딩 생성 (작은 배치 + 429 재시도)
                    texts = [doc.page_content for doc in docs]
                    embedding_vectors = embed_documents_with_retry(embeddings, texts)

                    # 임베딩 성공 후 기존 문서를 삭제해야 중간 실패 시 기존 검색 데이터가 보존됩니다.
                    delete_document_by_source(supabase_client, fname)

                    # RPC 함수로 안전한 삽입 (FLOAT[] → VECTOR(768))
                    for doc, embedding_vector in zip(docs, embedding_vectors):
                        supabase_client.rpc("insert_document_safe", {
                            "p_content": doc.page_content,
                            "p_metadata": doc.metadata,
                            "p_embedding_array": embedding_vector
                        }).execute()

                    st.success(f"[OK] {fname} 완료 ({len(docs)}개 청크)")
                    cnt += 1
                except Exception as insert_error:
                    st.error(f"[ERROR] {fname} 삽입 실패: {str(insert_error)[:100]}")
                    print(f"삽입 에러 - {fname}: {insert_error}")
                    failed += 1
            else:
                st.warning(f"[WARNING] {fname} - 유효한 청크 없음")
                skipped += 1

        except Exception as e:
            st.error(f"[ERROR] {fname} 실패: {str(e)[:100]}")
            print(f"상세 에러 - {fname}: {e}")
            failed += 1

    progress.empty()

    # 결과 요약
    if not force_update:
        st.info(f"""
        증분 동기화 완료
        - 색인 성공: {cnt}개
        - 삭제 처리: {deleted_count}개
        - 건너뜀: {skipped}개
        - 실패: {failed}개
        """)
    else:
        st.info(f"""
        전체 재색인 완료
        - 성공: {cnt}개
        - 건너뜀: {skipped}개
        - 실패: {failed}개
        - 전체: {len(files)}개
        """)

    return cnt + deleted_count

def search_similar_documents_with_retry(query, client, embeddings, top_k=5, threshold=0.5, max_retries=3, use_doc_level_ranking=True):
    """
    하이브리드 검색 함수 (Vector + Keyword)
    - 벡터 유사도 60% + 키워드 매칭 40%
    - 정확한 단어가 있으면 상위 노출 보장
    - 문서 레벨 랭킹으로 청크 수가 많은 문서의 불공정한 우위 방지
    """
    for attempt in range(max_retries):
        try:
            query_vector = embeddings.embed_query(query)

            # 문서 레벨 랭킹 검색 시도 (우선)
            if use_doc_level_ranking:
                try:
                    params = {
                        "query_embedding": query_vector,
                        "query_text": query,
                        "match_threshold": threshold,
                        "match_count": top_k,  # 상위 N개 문서
                        "keyword_weight": 0.4,
                        "chunks_per_doc": 2  # 각 문서당 2개 청크 반환
                    }
                    response = client.rpc("hybrid_search_documents_doc_ranked", params).execute()
                    use_doc_ranking = True
                    use_hybrid = True
                except Exception as e:
                    print(f"문서 레벨 랭킹 실패, 청크 레벨 랭킹으로 폴백: {e}")
                    use_doc_ranking = False
                    use_hybrid = False
            else:
                use_doc_ranking = False
                use_hybrid = False

            # 청크 레벨 하이브리드 검색 (폴백)
            if not use_doc_ranking:
                try:
                    params = {
                        "query_embedding": query_vector,
                        "query_text": query,
                        "match_threshold": threshold,
                        "match_count": top_k,
                        "keyword_weight": 0.4
                    }
                    response = client.rpc("hybrid_search_documents", params).execute()
                    use_hybrid = True
                except Exception as e:
                    print(f"하이브리드 검색 실패, 벡터 검색으로 폴백: {e}")
                    # 벡터 검색 (최종 폴백)
                    params = {
                        "query_embedding": query_vector,
                        "match_threshold": threshold,
                        "match_count": top_k
                    }
                    response = client.rpc("match_documents", params).execute()
                    use_hybrid = False

            docs = []
            infos = []

            for item in response.data:
                content = item.get("content", "")
                metadata = item.get("metadata", {})

                # 점수 추출
                if use_hybrid:
                    score = item.get("hybrid_score", item.get("similarity", 0.0))
                else:
                    score = item.get("similarity", 0.0)

                doc = Document(page_content=content, metadata=metadata)
                docs.append(doc)

                infos.append({
                    "content": content,
                    "filename": metadata.get("source", "Unknown"),
                    "section": metadata.get("section", "일반"),
                    "score": score
                })

            return docs, infos

        except Exception as e:
            if attempt == max_retries - 1:
                print(f"검색 실패 ({max_retries}회 시도): {e}")
                return [], []
            print(f"검색 재시도 {attempt + 1}/{max_retries}...")
            time.sleep(1 * (attempt + 1))

def search_similar_documents(query, client, embeddings, top_k=5, dynamic_threshold=True):
    """
    동적 임계값을 사용하는 개선된 검색 함수

    Args:
        query: 검색 쿼리
        client: Supabase 클라이언트
        embeddings: 임베딩 모델
        top_k: 반환할 문서 개수
        dynamic_threshold: True일 경우 동적 임계값 사용

    개선사항:
        - 임계값 대폭 낮춤: 0.5 → 0.3, 0.15 → 0.1
        - 더 많은 관련 문서 검색 가능
    """
    if dynamic_threshold:
        # 먼저 중간 임계값으로 검색 (0.5 → 0.3)
        docs_high, infos_high = search_similar_documents_with_retry(
            query, client, embeddings, top_k=top_k, threshold=0.3  # 🔧 0.5 → 0.3
        )

        # 결과가 충분하면 반환 (3개 → 5개로 상향)
        if len(docs_high) >= 5:  # 🔧 3 → 5
            return docs_high, infos_high

        # 결과 부족 시 낮은 임계값으로 재검색 (0.15 → 0.1)
        docs_low, infos_low = search_similar_documents_with_retry(
            query, client, embeddings, top_k=top_k, threshold=0.1  # 🔧 0.15 → 0.1
        )
        return docs_low, infos_low
    else:
        # 고정 임계값 사용 (0.15 → 0.1)
        return search_similar_documents_with_retry(
            query, client, embeddings, top_k=top_k, threshold=0.1  # 🔧 0.15 → 0.1
        )

def get_indexed_documents(client):
    """DB에 색인된 문서 목록 조회"""
    try:
        result = client.table("documents").select("metadata").execute()
        sources = set()
        for item in result.data:
            metadata = item.get("metadata", {})
            source = metadata.get("source", "Unknown")
            sources.add(source)
        return list(sources)
    except Exception as e:
        print(f"문서 목록 조회 실패: {e}")
        return []

def reset_database(client):
    """
    Supabase DB의 documents 테이블 전체를 삭제합니다. (강제 초기화)

    개선된 삭제 방법:
    1. 먼저 전체 문서 개수 확인
    2. 배치 단위로 삭제 (Supabase 제한 고려)
    3. 삭제 후 재확인
    """
    if client is None:
        print("Error: Supabase 클라이언트가 유효하지 않습니다.")
        return False

    try:
        # 1. 삭제 전 문서 개수 확인
        count_before = client.table("documents").select("id", count="exact").execute()
        total_docs = count_before.count if hasattr(count_before, 'count') else 0
        print(f"삭제 전 문서 개수: {total_docs}")

        if total_docs == 0:
            print("이미 빈 DB입니다.")
            return True

        # 2. 전체 삭제 실행 (neq 조건으로 모든 행 삭제)
        # Supabase는 id가 UUID이므로 gte('id', '00000000-0000-0000-0000-000000000000') 사용
        result = client.table("documents").delete().gte("id", "00000000-0000-0000-0000-000000000000").execute()

        print(f"삭제 실행 완료. 응답: {result}")

        # 3. 삭제 후 재확인
        count_after = client.table("documents").select("id", count="exact").execute()
        remaining_docs = count_after.count if hasattr(count_after, 'count') else 0
        print(f"삭제 후 남은 문서: {remaining_docs}")

        if remaining_docs == 0:
            print(f"[OK] DB 삭제 성공: {total_docs}개 문서 삭제됨")
            return True
        else:
            print(f"[WARNING] 완전 삭제 실패: {remaining_docs}개 문서 남음")
            return False

    except Exception as e:
        print(f"DB Reset Error: {e}")
        import traceback
        traceback.print_exc()
        return False
